"""
Document Operations Skill — SmartDome OS v6
=============================================
Generates office documents (XLSX, DOCX, PPTX, PDF) from agent data.
Files are created in /tmp, then uploaded to Google Drive via drive_tool.
This is what gives agents the ability to produce real deliverables.
"""
import json
import io
import logging
import tempfile
from pathlib import Path
from datetime import datetime
from typing import Dict, Any, List, Optional

logger = logging.getLogger("SmartDome-DocOps")

# --- Path Resolution ---
_SCRIPT_DIR = Path(__file__).resolve().parent
_SERVER_ROOT = _SCRIPT_DIR.parent.parent

if "Personal assistant" in str(_SERVER_ROOT):
    _DATA_DIR = _SERVER_ROOT / "data"
else:
    _DATA_DIR = Path("/tmp/data")

_OUTPUT_DIR = Path(tempfile.gettempdir()) / "smartdome_docs"
_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


# =====================================================================
# EXCEL (XLSX) GENERATION
# =====================================================================

def create_spreadsheet(
    title: str,
    headers: str,
    rows: str,
    sheet_name: str = "Sheet1",
    folder_path: str = "Shared",
    agent_id: str = "cfo"
) -> Dict[str, Any]:
    """
    Creates an Excel spreadsheet (.xlsx) with data and saves it.
    The CFO uses this for budgets, invoices, and financial reports.
    The CEO uses this for operational reports.

    Args:
        title: Document title / filename (without .xlsx extension).
        headers: Comma-separated column headers (e.g., 'Date,Description,Amount,Category').
        rows: JSON array of arrays, each inner array is a row (e.g., '[["2026-01-15","Vicat payment",5000,"Materials"]]').
        sheet_name: Name for the worksheet tab. Default 'Sheet1'.
        folder_path: Drive folder to save to (e.g., 'Finance/Invoices'). Default 'Shared'.
        agent_id: The agent creating the document.

    Returns:
        dict: File path, size, and status. Ready for Drive upload.
    """
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = sheet_name

        # Parse inputs
        header_list = [h.strip() for h in headers.split(",")]
        try:
            row_data = json.loads(rows) if isinstance(rows, str) else rows
        except json.JSONDecodeError:
            return {"success": False, "error": "Invalid rows JSON. Expected array of arrays."}

        # Style definitions
        header_fill = PatternFill(start_color="1A56DB", end_color="1A56DB", fill_type="solid")
        header_font = Font(name="Arial", bold=True, color="FFFFFF", size=11)
        border = Border(
            left=Side(style="thin", color="CBD5E1"),
            right=Side(style="thin", color="CBD5E1"),
            top=Side(style="thin", color="CBD5E1"),
            bottom=Side(style="thin", color="CBD5E1"),
        )

        # Write headers
        for col, header in enumerate(header_list, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center")
            cell.border = border

        # Write data rows
        for row_idx, row in enumerate(row_data, 2):
            for col_idx, value in enumerate(row, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=value)
                cell.border = border
                cell.font = Font(name="Arial", size=10)
                # Auto-detect numbers for right alignment
                if isinstance(value, (int, float)):
                    cell.alignment = Alignment(horizontal="right")
                    cell.number_format = '#,##0.00'

        # Auto-width columns
        for col_idx, header in enumerate(header_list, 1):
            max_len = len(header)
            for row_idx in range(2, len(row_data) + 2):
                cell_val = str(ws.cell(row=row_idx, column=col_idx).value or "")
                max_len = max(max_len, len(cell_val))
            ws.column_dimensions[openpyxl.utils.get_column_letter(col_idx)].width = min(max_len + 4, 40)

        # Add auto-filter
        ws.auto_filter.ref = ws.dimensions

        # Freeze header row
        ws.freeze_panes = "A2"

        # Save
        filename = f"{title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.xlsx"
        filepath = _OUTPUT_DIR / filename
        wb.save(str(filepath))

        file_size = filepath.stat().st_size
        logger.info(f"Created spreadsheet: {filename} ({file_size} bytes) by {agent_id}")

        return {
            "success": True,
            "message": f"Spreadsheet '{filename}' created ({len(row_data)} rows)",
            "filename": filename,
            "filepath": str(filepath),
            "size_bytes": file_size,
            "drive_folder": folder_path,
            "rows": len(row_data),
            "columns": len(header_list),
            "ready_for_upload": True
        }
    except ImportError:
        return {"success": False, "error": "openpyxl not installed. pip install openpyxl"}
    except Exception as e:
        return {"success": False, "error": str(e)}


# =====================================================================
# WORD DOCUMENT (DOCX) GENERATION
# =====================================================================

def create_document(
    title: str,
    content_sections: str,
    author: str = "SmartDome OS",
    folder_path: str = "Shared",
    agent_id: str = "ceo"
) -> Dict[str, Any]:
    """
    Creates a Word document (.docx) with formatted sections.
    Used by CEO for reports, CLO for contracts, CMO for marketing docs.

    Args:
        title: Document title (used as filename and heading).
        content_sections: JSON array of objects with 'heading' and 'body' keys.
            Example: '[{"heading":"Summary","body":"Q4 was strong..."},{"heading":"Next Steps","body":"We plan to..."}]'
        author: Document author name. Default 'SmartDome OS'.
        folder_path: Drive folder to save to. Default 'Shared'.
        agent_id: The agent creating the document.

    Returns:
        dict: File path, size, and status.
    """
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        # Set default font
        style = doc.styles["Normal"]
        font = style.font
        font.name = "Arial"
        font.size = Pt(11)

        # Title
        title_para = doc.add_heading(title, level=0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Metadata line
        meta = doc.add_paragraph()
        meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = meta.add_run(f"Generated by {author} | {datetime.now().strftime('%B %d, %Y')}")
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(148, 163, 184)

        doc.add_paragraph()  # Spacer

        # Parse sections
        try:
            sections = json.loads(content_sections) if isinstance(content_sections, str) else content_sections
        except json.JSONDecodeError:
            return {"success": False, "error": "Invalid content_sections JSON."}

        for section in sections:
            heading = section.get("heading", "")
            body = section.get("body", "")

            if heading:
                doc.add_heading(heading, level=1)
            if body:
                # Handle paragraphs separated by newlines
                for paragraph_text in body.split("\n"):
                    if paragraph_text.strip():
                        doc.add_paragraph(paragraph_text.strip())

        # Footer with confidentiality
        section = doc.sections[0]
        footer = section.footer
        footer_para = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
        footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = footer_para.add_run("CONFIDENTIAL — SmartDome OS")
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(148, 163, 184)

        # Save
        filename = f"{title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.docx"
        filepath = _OUTPUT_DIR / filename
        doc.save(str(filepath))

        file_size = filepath.stat().st_size
        logger.info(f"Created document: {filename} ({file_size} bytes) by {agent_id}")

        return {
            "success": True,
            "message": f"Document '{filename}' created ({len(sections)} sections)",
            "filename": filename,
            "filepath": str(filepath),
            "size_bytes": file_size,
            "drive_folder": folder_path,
            "sections": len(sections),
            "ready_for_upload": True
        }
    except ImportError:
        return {"success": False, "error": "python-docx not installed. pip install python-docx"}
    except Exception as e:
        return {"success": False, "error": str(e)}


# =====================================================================
# PRESENTATION (PPTX) GENERATION
# =====================================================================

def create_presentation(
    title: str,
    slides: str,
    author: str = "SmartDome OS",
    folder_path: str = "Shared",
    agent_id: str = "cmo"
) -> Dict[str, Any]:
    """
    Creates a PowerPoint presentation (.pptx) with title and content slides.
    Used by CMO for marketing decks, CEO for strategy presentations.

    Args:
        title: Presentation title (filename and title slide).
        slides: JSON array of objects with 'title' and 'content' keys.
            Example: '[{"title":"Market Overview","content":"SmartDome targets the premium housing segment..."}]'
        author: Presentation author. Default 'SmartDome OS'.
        folder_path: Drive folder to save to. Default 'Shared'.
        agent_id: The agent creating the presentation.

    Returns:
        dict: File path, size, slide count, and status.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)

        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        subtitle = slide.placeholders[1]
        subtitle.text = f"{author} | {datetime.now().strftime('%B %Y')}"

        # Parse slides
        try:
            slide_data = json.loads(slides) if isinstance(slides, str) else slides
        except json.JSONDecodeError:
            return {"success": False, "error": "Invalid slides JSON."}

        for s in slide_data:
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = s.get("title", "")

            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()

            content = s.get("content", "")
            for i, line in enumerate(content.split("\n")):
                if i == 0:
                    tf.text = line.strip()
                else:
                    p = tf.add_paragraph()
                    p.text = line.strip()

        # Save
        filename = f"{title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.pptx"
        filepath = _OUTPUT_DIR / filename
        prs.save(str(filepath))

        file_size = filepath.stat().st_size
        logger.info(f"Created presentation: {filename} ({file_size} bytes) by {agent_id}")

        return {
            "success": True,
            "message": f"Presentation '{filename}' created ({len(slide_data) + 1} slides)",
            "filename": filename,
            "filepath": str(filepath),
            "size_bytes": file_size,
            "drive_folder": folder_path,
            "slide_count": len(slide_data) + 1,
            "ready_for_upload": True
        }
    except ImportError:
        return {"success": False, "error": "python-pptx not installed. pip install python-pptx"}
    except Exception as e:
        return {"success": False, "error": str(e)}


# =====================================================================
# INVOICE MANAGEMENT (CFO SPECIAL)
# =====================================================================

def record_invoice(
    vendor: str,
    amount: float,
    currency: str = "BGN",
    invoice_number: str = "",
    description: str = "",
    category: str = "general",
    due_date: str = ""
) -> Dict[str, Any]:
    """
    Records an invoice in the SmartDome invoice ledger and creates a tracking entry.
    The CFO agent uses this to manage all incoming and outgoing invoices.

    Args:
        vendor: Company or person name on the invoice.
        amount: Invoice amount (positive = payable, negative = receivable).
        currency: Currency code. Default 'BGN'.
        invoice_number: Official invoice number/ID.
        description: What the invoice is for.
        category: Category - 'materials', 'services', 'saas', 'legal', 'travel', 'general'.
        due_date: Payment due date (YYYY-MM-DD format). Empty = immediate.

    Returns:
        dict: Invoice record with tracking ID and storage path.
    """
    try:
        _DATA_DIR.mkdir(parents=True, exist_ok=True)
        invoice_file = _DATA_DIR / "smartdome" / "invoice_ledger.json"
        invoice_file.parent.mkdir(parents=True, exist_ok=True)

        data = {"invoices": [], "summary": {"total_payable": 0, "total_receivable": 0, "count": 0}}
        if invoice_file.exists():
            with open(invoice_file, "r", encoding="utf-8") as f:
                data = json.load(f)

        invoice = {
            "id": f"INV-{datetime.now().strftime('%Y%m%d')}-{len(data['invoices']) + 1:04d}",
            "vendor": vendor,
            "amount": amount,
            "currency": currency,
            "invoice_number": invoice_number or "N/A",
            "description": description,
            "category": category.lower(),
            "due_date": due_date or "immediate",
            "status": "recorded",
            "recorded_at": datetime.now().isoformat(),
            "recorded_by": "CFO",
            "drive_path": f"Finance/Invoices/{datetime.now().strftime('%Y-%m')}"
        }

        data["invoices"].insert(0, invoice)
        data["summary"]["count"] = len(data["invoices"])
        data["summary"]["total_payable"] = sum(
            i["amount"] for i in data["invoices"] if i["amount"] > 0 and i["status"] != "paid"
        )
        data["summary"]["total_receivable"] = sum(
            abs(i["amount"]) for i in data["invoices"] if i["amount"] < 0 and i["status"] != "received"
        )

        with open(invoice_file, "w", encoding="utf-8") as f:
            json.dump(data, f, indent=2, ensure_ascii=False)

        logger.info(f"Invoice recorded: {invoice['id']} — {vendor} {amount} {currency}")

        return {
            "success": True,
            "message": f"Invoice {invoice['id']} recorded: {vendor} — {amount} {currency}",
            "invoice": invoice,
            "summary": data["summary"]
        }
    except Exception as e:
        return {"success": False, "error": str(e)}


def get_invoice_summary(status: str = "all", period: str = "all") -> Dict[str, Any]:
    """
    Returns invoice ledger summary with filters.

    Args:
        status: Filter by status - 'all', 'recorded', 'paid', 'overdue'.
        period: Time period - 'all', 'month', 'week'.

    Returns:
        dict: Invoice summary with totals, breakdown by vendor, and recent entries.
    """
    try:
        invoice_file = _DATA_DIR / "smartdome" / "invoice_ledger.json"
        if not invoice_file.exists():
            return {"success": True, "message": "No invoices recorded yet.", "invoices": [], "totals": {}}

        with open(invoice_file, "r", encoding="utf-8") as f:
            data = json.load(f)

        invoices = data.get("invoices", [])

        # Filter by status
        if status != "all":
            invoices = [i for i in invoices if i.get("status") == status]

        # Filter by period
        if period == "month":
            cutoff = datetime.now().replace(day=1).isoformat()
            invoices = [i for i in invoices if i.get("recorded_at", "") >= cutoff]
        elif period == "week":
            from datetime import timedelta
            cutoff = (datetime.now() - timedelta(days=7)).isoformat()
            invoices = [i for i in invoices if i.get("recorded_at", "") >= cutoff]

        # Group by vendor
        by_vendor = {}
        for i in invoices:
            v = i.get("vendor", "Unknown")
            by_vendor[v] = by_vendor.get(v, 0) + i.get("amount", 0)

        total_payable = sum(i["amount"] for i in invoices if i["amount"] > 0)
        total_receivable = sum(abs(i["amount"]) for i in invoices if i["amount"] < 0)

        return {
            "success": True,
            "period": period,
            "status_filter": status,
            "total_payable": total_payable,
            "total_receivable": total_receivable,
            "net": total_receivable - total_payable,
            "by_vendor": by_vendor,
            "count": len(invoices),
            "recent": invoices[:10]
        }
    except Exception as e:
        return {"success": False, "error": str(e)}


def generate_invoice_report(period: str = "month", format: str = "xlsx") -> Dict[str, Any]:
    """
    Generates a formatted invoice report as an Excel spreadsheet.
    Combines invoice ledger data into a professional financial document.

    Args:
        period: Time period for the report - 'all', 'month', 'quarter', 'week'.
        format: Output format - 'xlsx' (default).

    Returns:
        dict: Generated report filepath and summary stats.
    """
    try:
        # Get invoice data
        summary = get_invoice_summary(period=period)
        if not summary.get("success"):
            return summary

        invoices = summary.get("recent", [])
        if not invoices:
            return {"success": False, "error": f"No invoices found for period: {period}"}

        # Build spreadsheet data
        headers = "Date,Invoice #,Vendor,Description,Category,Amount,Currency,Status,Due Date"
        rows = []
        for inv in invoices:
            rows.append([
                inv.get("recorded_at", "")[:10],
                inv.get("invoice_number", "N/A"),
                inv.get("vendor", ""),
                inv.get("description", ""),
                inv.get("category", ""),
                inv.get("amount", 0),
                inv.get("currency", "BGN"),
                inv.get("status", ""),
                inv.get("due_date", "")
            ])

        result = create_spreadsheet(
            title=f"Invoice_Report_{period}",
            headers=headers,
            rows=json.dumps(rows),
            sheet_name="Invoices",
            folder_path="Finance/Reports",
            agent_id="cfo"
        )

        if result.get("success"):
            result["summary"] = {
                "total_payable": summary["total_payable"],
                "total_receivable": summary["total_receivable"],
                "net": summary["net"],
                "invoice_count": summary["count"]
            }

        return result
    except Exception as e:
        return {"success": False, "error": str(e)}
